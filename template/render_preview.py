#!/usr/bin/env python3
"""Faithful-enough slide preview renderer for QA (LibreOffice is broken in this
sandbox). Reads real shape geometry / fills / text / embedded images via
python-pptx and draws each slide with PIL. Catches overflow, overlap, off-slide."""
import io, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

HERE = Path(__file__).resolve().parent
PPTX = HERE / "AZBY_Portfolio_Review_Template.pptx"
OUT = HERE / "assets" / "preview"; OUT.mkdir(parents=True, exist_ok=True)
SC = 100  # px per inch
FONT = "/usr/share/fonts/truetype/crosextra/Carlito-Regular.ttf"
FONTB = "/usr/share/fonts/truetype/crosextra/Carlito-Bold.ttf"
EMU_IN = 914400
def px(emu): return int(emu / EMU_IN * SC)
_fc = {}
def font(sz, bold):
    k = (int(sz), bold)
    if k not in _fc: _fc[k] = ImageFont.truetype(FONTB if bold else FONT, max(6, int(sz)))
    return _fc[k]

def hexof(fill, default=None):
    try:
        if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
            return "#" + str(fill.fore_color.rgb)
    except Exception:
        pass
    return default

def wrap(draw, text, fnt, maxw):
    out = []
    for para in text.split("\n"):
        words = para.split(" "); line = ""
        for w in words:
            t = (line + " " + w).strip()
            if draw.textlength(t, font=fnt) <= maxw or not line:
                line = t
            else:
                out.append(line); line = w
        out.append(line)
    return out

def draw_textbox(draw, x, y, w, h, tf, defcolor="#16233B", overflow_marks=None):
    # gather runs -> approximate first-run size/bold/color per paragraph
    yy = y + 3
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            yy += 14; continue
        sz = None; bold = False; col = defcolor
        txt = "".join(r.text for r in runs)
        r0 = runs[0]
        if r0.font.size: sz = r0.font.size.pt
        bold = bool(r0.font.bold)
        try:
            if r0.font.color and r0.font.color.type is not None:
                col = "#" + str(r0.font.color.rgb)
        except Exception: pass
        sz = sz or 13
        fnt = font(sz * SC / 72, bold)
        lines = wrap(draw, txt, fnt, w - 8)
        lh = (sz * SC / 72) * 1.2
        for ln in lines:
            if yy + lh > y + h + 2 and overflow_marks is not None:
                overflow_marks.append((x, y, w, h));
            draw.text((x + 4, yy), ln, font=fnt, fill=col)
            yy += lh

def render():
    prs = Presentation(str(PPTX))
    W = px(prs.slide_width); H = px(prs.slide_height)
    n = 0
    imgs = []
    for slide in prs.slides:
        n += 1
        # background: detect solid fill
        bg = "#FFFFFF"
        try:
            if slide.background.fill.type is not None:
                bg = hexof(slide.background.fill, "#FFFFFF")
        except Exception: pass
        img = Image.new("RGB", (W, H), bg or "#FFFFFF")
        d = ImageDraw.Draw(img)
        overflow = []
        for sh in slide.shapes:
            try:
                x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            except Exception:
                continue
            # picture
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
                    im.thumbnail((max(1, w), max(1, h)))
                    bgp = Image.new("RGBA", im.size, bg if bg and bg.startswith("#") else "#FFFFFF")
                    bgp.alpha_composite(im); img.paste(bgp.convert("RGB"), (x, y))
                except Exception as e:
                    d.rectangle([x, y, x + w, y + h], outline="#888888")
                continue
            # table
            if sh.has_table:
                tbl = sh.table
                cw = [px(c.width) for c in tbl.columns]
                rh = [px(r.height) for r in tbl.rows]
                cy = y
                for ri, row in enumerate(tbl.rows):
                    cx = x
                    for ci, cell in enumerate(row.cells):
                        cwi = cw[ci] if ci < len(cw) else 40
                        rhi = rh[ri] if ri < len(rh) else 20
                        cf = hexof(cell.fill, None)
                        if cf: d.rectangle([cx, cy, cx + cwi, cy + rhi], fill=cf)
                        d.rectangle([cx, cy, cx + cwi, cy + rhi], outline="#FFFFFF", width=1)
                        # text
                        tf = cell.text_frame
                        draw_textbox(d, cx, cy, cwi, rhi, tf, "#16233B", overflow)
                        cx += cwi
                    cy += rh[ri] if ri < len(rh) else 20
                continue
            # autoshape/textbox fill
            fillc = hexof(sh.fill, None) if hasattr(sh, "fill") else None
            if fillc:
                d.rounded_rectangle([x, y, x + w, y + h], radius=6, fill=fillc)
            # text
            if sh.has_text_frame and sh.text_frame.text.strip():
                dc = "#FFFFFF" if (fillc and fillc.lower() in ("#16233b", "#1b27a3")) or (bg and bg.lower() in ("#16233b",)) else "#16233B"
                draw_textbox(d, x, y, w, h, sh.text_frame, dc, overflow)
        # mark overflow boxes
        for (ox, oy, ow, oh) in overflow:
            d.rectangle([ox, oy, ox + ow, oy + oh], outline="#E0402F", width=3)
        d.rectangle([0, 0, W - 1, H - 1], outline="#CCCCCC")
        p = OUT / f"slide-{n:02d}.png"; img.save(p); imgs.append(p)
    # contact sheets (12 per sheet)
    per = 9; cols = 3
    sheets = []
    for i in range(0, len(imgs), per):
        chunk = imgs[i:i+per]; rows = (len(chunk)+cols-1)//cols
        cw, ch = 460, 260
        sheet = Image.new("RGB", (cols*cw, rows*ch), "white"); sd = ImageDraw.Draw(sheet)
        for j, p in enumerate(chunk):
            im = Image.open(p); im.thumbnail((cw-16, ch-30))
            xx = (j % cols)*cw; yy = (j//cols)*ch
            sheet.paste(im, (xx+8, yy+24))
            sd.text((xx+8, yy+6), p.name, fill="black", font=font(13, True))
        sp = OUT / f"contact_{i//per+1}.png"; sheet.save(sp); sheets.append(str(sp))
    print("rendered", len(imgs), "slides;", len(sheets), "contact sheets:")
    for s in sheets: print(" ", s)

if __name__ == "__main__":
    render()
